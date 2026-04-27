from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Palette ───────────────────────────────────────────────────────────────────
NAVY      = RGBColor(0x0A, 0x1F, 0x44)   # slide background
RED       = RGBColor(0xE0, 0x2B, 0x2B)   # accent / ESI red
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT     = RGBColor(0xE8, 0xEF, 0xF8)   # light text / table header bg
GOLD      = RGBColor(0xF5, 0xA6, 0x23)   # highlight accent
GRAY_BG   = RGBColor(0x12, 0x2B, 0x55)   # card / table row bg
MID_BLUE  = RGBColor(0x1A, 0x3A, 0x6B)   # alternate row
MONO_BG   = RGBColor(0x0D, 0x1B, 0x36)   # code block bg

W = Inches(13.33)   # widescreen width
H = Inches(7.5)     # widescreen height

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank layout


# ── Helpers ───────────────────────────────────────────────────────────────────

def add_slide():
    return prs.slides.add_slide(BLANK)

def fill_bg(slide, color=NAVY):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(
        1, left, top, width, height)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def add_label(slide, text, left, top, width, height,
              font_size=18, bold=False, color=WHITE,
              align=PP_ALIGN.LEFT, italic=False):
    return add_textbox(slide, text, left, top, width, height,
                       font_size=font_size, bold=bold, color=color,
                       align=align, wrap=True, italic=italic)

def accent_bar(slide):
    """Red bar on the left edge."""
    add_rect(slide, Inches(0), Inches(0), Inches(0.18), H, RED)

def slide_header(slide, title, subtitle=None):
    accent_bar(slide)
    add_label(slide, title,
              Inches(0.35), Inches(0.28), Inches(12.6), Inches(0.8),
              font_size=32, bold=True, color=WHITE)
    if subtitle:
        add_label(slide, subtitle,
                  Inches(0.35), Inches(1.0), Inches(12.6), Inches(0.5),
                  font_size=16, color=GOLD)

def footer(slide, text="BUAN 6v99.s01  ·  UT Dallas  ·  April 2026"):
    add_rect(slide, Inches(0), Inches(7.15), W, Inches(0.35), GRAY_BG)
    add_label(slide, text,
              Inches(0.35), Inches(7.17), Inches(12.6), Inches(0.3),
              font_size=9, color=LIGHT, align=PP_ALIGN.LEFT)

def bullets(slide, items, left, top, width, height,
            font_size=16, color=WHITE, indent=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p = tf.add_paragraph() if not first else tf.paragraphs[0]
        first = False
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = ("    " if indent else "• ") + item
        run.font.size = Pt(font_size)
        run.font.color.rgb = color

def simple_table(slide, headers, rows,
                 left, top, col_widths, row_height=Inches(0.42),
                 header_bg=RED, row_bg=GRAY_BG, alt_bg=MID_BLUE,
                 font_size=13):
    """Draw a table manually as rectangles + text (no pptx table object)."""
    n_cols = len(headers)
    # header row
    x = left
    for i, (h, cw) in enumerate(zip(headers, col_widths)):
        add_rect(slide, x, top, cw, row_height, header_bg)
        add_label(slide, h, x + Inches(0.08), top + Inches(0.06),
                  cw - Inches(0.12), row_height - Inches(0.08),
                  font_size=font_size, bold=True, color=WHITE)
        x += cw
    # data rows
    for r_idx, row in enumerate(rows):
        y = top + row_height * (r_idx + 1)
        bg = row_bg if r_idx % 2 == 0 else alt_bg
        x = left
        for c_idx, (cell, cw) in enumerate(zip(row, col_widths)):
            add_rect(slide, x, y, cw, row_height, bg)
            add_label(slide, str(cell),
                      x + Inches(0.08), y + Inches(0.06),
                      cw - Inches(0.12), row_height - Inches(0.08),
                      font_size=font_size, color=WHITE)
            x += cw


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
fill_bg(s)

# Big red bar top
add_rect(s, Inches(0), Inches(0), W, Inches(0.12), RED)
# Left accent
add_rect(s, Inches(0), Inches(0), Inches(0.18), H, RED)

# Logo/brand pill
add_rect(s, Inches(0.45), Inches(1.2), Inches(2.6), Inches(0.55), RED)
add_label(s, "TriageIQ", Inches(0.5), Inches(1.22), Inches(2.5), Inches(0.5),
          font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Title
add_label(s, "AI-Assisted Emergency Triage",
          Inches(0.45), Inches(1.95), Inches(12), Inches(1.1),
          font_size=44, bold=True, color=WHITE)

# Subtitle
add_label(s, "Clinical Decision Support for Emergency Nurses",
          Inches(0.45), Inches(3.1), Inches(12), Inches(0.6),
          font_size=22, color=GOLD)

# Tagline
add_label(s, '"Not replacing the nurse.  Backing them up."',
          Inches(0.45), Inches(3.85), Inches(10), Inches(0.5),
          font_size=16, italic=True, color=LIGHT)

# Divider
add_rect(s, Inches(0.45), Inches(4.55), Inches(6), Inches(0.04), RED)

# Footer info
add_label(s, "BUAN 6v99.s01  ·  UT Dallas  ·  April 2026",
          Inches(0.45), Inches(4.75), Inches(8), Inches(0.4),
          font_size=13, color=LIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — The Problem
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
fill_bg(s)
slide_header(s, "The Problem",
             "ER Triage Is High-Stakes and Time-Pressured")
footer(s)

pts = [
    "Emergency departments see millions of patients annually — nurses must assign an ESI score within minutes of arrival",
    "Cognitive overload: nurses simultaneously assess vitals, chief complaint, history, and resource availability under pressure",
    "Inconsistent triage scoring increases risk: under-triaging delays critical care; over-triaging wastes scarce resources",
    "Protocol knowledge is vast — no single nurse can recall every guideline in real time",
    "Consequence of error: delayed intervention for STEMI, stroke, sepsis, or respiratory failure can cost lives",
]
bullets(s, pts, Inches(0.45), Inches(1.65), Inches(12.5), Inches(5.2), font_size=17)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — The Gap
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
fill_bg(s)
slide_header(s, "The Gap", "Existing Tools Don't Go Far Enough")
footer(s)

pts = [
    "EHR systems record data but don't synthesize it into a recommendation",
    "Protocol lookup tools exist but require the nurse to know what to search for",
    "No widely adopted tool brings vitals analysis, symptom classification, protocol matching, and resource allocation together in one workflow",
    "Nurses are left to mentally integrate all inputs under time pressure — alone",
]
bullets(s, pts, Inches(0.45), Inches(1.65), Inches(12.5), Inches(3.5), font_size=18)

# Visual: nurse-as-hub diagram (text-based)
add_rect(s, Inches(1.5), Inches(4.5), Inches(10.3), Inches(2.3), GRAY_BG)
add_label(s, "Vitals Monitor", Inches(1.7), Inches(4.65), Inches(2.2), Inches(0.4),
          font_size=13, color=GOLD, bold=True)
add_label(s, "EHR System", Inches(4.3), Inches(4.65), Inches(2.2), Inches(0.4),
          font_size=13, color=GOLD, bold=True)
add_label(s, "Protocol Binder", Inches(6.9), Inches(4.65), Inches(2.4), Inches(0.4),
          font_size=13, color=GOLD, bold=True)
add_label(s, "Bed Board", Inches(9.5), Inches(4.65), Inches(1.8), Inches(0.4),
          font_size=13, color=GOLD, bold=True)
add_label(s, "↓                        ↓                        ↓                       ↓",
          Inches(1.7), Inches(5.05), Inches(9.8), Inches(0.35), font_size=14, color=LIGHT)
add_rect(s, Inches(5.2), Inches(5.45), Inches(2.9), Inches(0.55), RED)
add_label(s, "NURSE  (mental integration)",
          Inches(5.25), Inches(5.5), Inches(2.8), Inches(0.45),
          font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_label(s, "Single point of failure under time pressure",
          Inches(3.5), Inches(6.1), Inches(6.3), Inches(0.4),
          font_size=12, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — The Solution
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
fill_bg(s)
slide_header(s, "The Solution", "TriageIQ — An AI Triage Team in Your Browser")
footer(s)

pts = [
    "Simulates a team of specialist AI agents analyzing the patient in parallel",
    "Input: patient vitals, chief complaint, and medical history (EHR or free-text)",
    "Output in seconds: ESI score (1–5), prioritized action checklist, matched clinical protocols, and bed recommendation",
    "The nurse makes the final call — all recommendations are auditable and overridable with a documented reason",
]
bullets(s, pts, Inches(0.45), Inches(1.65), Inches(8.2), Inches(3.8), font_size=17)

# Output card
add_rect(s, Inches(9.0), Inches(1.5), Inches(4.0), Inches(5.3), GRAY_BG)
add_rect(s, Inches(9.0), Inches(1.5), Inches(4.0), Inches(0.55), RED)
add_label(s, "OUTPUT", Inches(9.0), Inches(1.52), Inches(4.0), Inches(0.5),
          font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
output_items = [
    "ESI Score  1–5",
    "Action Checklist",
    "Matched Protocol",
    "Bed Recommendation",
    "Nurse Override Panel",
    "Ask AI Chat",
]
for i, item in enumerate(output_items):
    add_rect(s, Inches(9.15), Inches(2.2) + Inches(0.68) * i,
             Inches(3.7), Inches(0.55), MID_BLUE)
    add_label(s, item,
              Inches(9.25), Inches(2.24) + Inches(0.68) * i,
              Inches(3.5), Inches(0.45),
              font_size=14, color=WHITE)

add_label(s, '"Not replacing the nurse.  Backing them up."',
          Inches(0.45), Inches(5.8), Inches(8.2), Inches(0.5),
          font_size=15, italic=True, color=GOLD)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Demo Preview
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
fill_bg(s)
slide_header(s, "Demo Preview", "What the Nurse Sees")
footer(s)

# Mock UI wireframe
# Patient list panel
add_rect(s, Inches(0.35), Inches(1.5), Inches(3.2), Inches(5.3), GRAY_BG)
add_rect(s, Inches(0.35), Inches(1.5), Inches(3.2), Inches(0.45), MID_BLUE)
add_label(s, "Patient Queue", Inches(0.45), Inches(1.52), Inches(3.0), Inches(0.4),
          font_size=12, bold=True, color=GOLD)
patient_cards = [
    ("PT-001", "STEMI", "ESI 1", RED),
    ("PT-002", "Subarachnoid", "ESI 1", RED),
    ("PT-003", "Ankle Fx", "ESI 3", RGBColor(0xF5,0xA6,0x23)),
    ("PT-004", "Peds Resp", "ESI 2", RGBColor(0xE0,0x60,0x00)),
    ("PT-005", "Elderly", "ESI 3", RGBColor(0xF5,0xA6,0x23)),
    ("PT-006", "Sepsis", "ESI 2", RGBColor(0xE0,0x60,0x00)),
]
for i, (pid, cc, esi, c) in enumerate(patient_cards):
    y = Inches(2.1) + Inches(0.73) * i
    add_rect(s, Inches(0.45), y, Inches(3.0), Inches(0.6), MID_BLUE)
    add_rect(s, Inches(0.45), y, Inches(0.08), Inches(0.6), c)
    add_label(s, f"{pid}  {cc}", Inches(0.6), y + Inches(0.04), Inches(2.0), Inches(0.3),
              font_size=11, color=WHITE, bold=True)
    add_label(s, esi, Inches(2.7), y + Inches(0.18), Inches(0.7), Inches(0.3),
              font_size=10, color=c, bold=True)

# Agent grid
add_rect(s, Inches(3.7), Inches(1.5), Inches(9.3), Inches(2.4), GRAY_BG)
add_rect(s, Inches(3.7), Inches(1.5), Inches(9.3), Inches(0.45), MID_BLUE)
add_label(s, "Analysis Progress", Inches(3.8), Inches(1.52), Inches(5.0), Inches(0.4),
          font_size=12, bold=True, color=GOLD)
agents = [
    ("Vitals", "✓ 1.2s", GREEN := RGBColor(0x22,0xC5,0x5E)),
    ("Symptoms", "✓ 1.8s", GREEN),
    ("Protocol", "⟳ running", RGBColor(0xF5,0xA6,0x23)),
    ("Beds", "⟳ running", RGBColor(0xF5,0xA6,0x23)),
    ("Synthesizer", "idle", LIGHT),
    ("Chat", "idle", LIGHT),
]
for i, (name, status, c) in enumerate(agents):
    col = i % 3
    row = i // 3
    x = Inches(3.8) + Inches(3.0) * col
    y = Inches(2.1) + Inches(0.8) * row
    add_rect(s, x, y, Inches(2.8), Inches(0.65), MID_BLUE)
    add_label(s, name, x + Inches(0.1), y + Inches(0.05), Inches(2.0), Inches(0.3),
              font_size=12, bold=True, color=WHITE)
    add_label(s, status, x + Inches(0.1), y + Inches(0.35), Inches(2.5), Inches(0.25),
              font_size=10, color=c)

# Report tabs
add_rect(s, Inches(3.7), Inches(4.1), Inches(9.3), Inches(0.45), MID_BLUE)
tabs = ["Summary", "Vitals", "Symptoms", "Protocol", "Resources", "Ask AI"]
for i, tab in enumerate(tabs):
    x = Inches(3.75) + Inches(1.52) * i
    bg = RED if i == 0 else GRAY_BG
    add_rect(s, x, Inches(4.12), Inches(1.45), Inches(0.4), bg)
    add_label(s, tab, x, Inches(4.14), Inches(1.45), Inches(0.36),
              font_size=10, bold=(i == 0), color=WHITE, align=PP_ALIGN.CENTER)

# ESI Banner
add_rect(s, Inches(3.7), Inches(4.65), Inches(9.3), Inches(0.7), RED)
add_label(s, "ESI 1  —  IMMEDIATE",
          Inches(3.8), Inches(4.7), Inches(9.1), Inches(0.6),
          font_size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_rect(s, Inches(3.7), Inches(5.5), Inches(9.3), Inches(1.3), GRAY_BG)
add_label(s, "Recommended Actions: Activate cath lab · 12-lead ECG · Aspirin 325 mg · IV access x2 · Cardiology consult",
          Inches(3.8), Inches(5.6), Inches(9.1), Inches(1.1),
          font_size=12, color=WHITE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Tech Stack
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
fill_bg(s)
slide_header(s, "Tech Stack", "Built on Modern AI and Web Infrastructure")
footer(s)

headers = ["Layer", "Technology"]
rows = [
    ["Frontend",          "React 18 + Vite 5  (JSX, custom CSS)"],
    ["Backend",           "Python 3.13  ·  FastAPI  ·  Uvicorn"],
    ["AI Agents",         "Anthropic Claude API  (Tool Use)"],
    ["Specialist Model",  "Claude Haiku 4.5 — fast, parallel agents"],
    ["Synthesis Model",   "Claude Sonnet 4.6 — final scoring and chat"],
    ["Streaming",         "Server-Sent Events (SSE)"],
    ["Concurrency",       "Python ThreadPoolExecutor"],
    ["Validation",        "Pydantic v2"],
    ["Data",              "In-memory mock (6 patients, bed state)"],
]
col_widths = [Inches(3.2), Inches(9.3)]
simple_table(s, headers, rows,
             Inches(0.35), Inches(1.6),
             col_widths, row_height=Inches(0.47), font_size=14)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Two-Tier Model Strategy
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
fill_bg(s)
slide_header(s, "Two-Tier Model Strategy", "Right Model for the Right Job")
footer(s)

# Left column — Haiku
add_rect(s, Inches(0.35), Inches(1.55), Inches(5.8), Inches(4.9), GRAY_BG)
add_rect(s, Inches(0.35), Inches(1.55), Inches(5.8), Inches(0.55), MID_BLUE)
add_label(s, "Claude Haiku  —  Fast + Cost-Efficient",
          Inches(0.45), Inches(1.58), Inches(5.6), Inches(0.48),
          font_size=14, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
haiku_items = ["Vitals Analyzer", "Symptom Classifier",
               "Protocol Matcher", "Bed Allocator", "Runs in PARALLEL"]
for i, item in enumerate(haiku_items):
    y = Inches(2.25) + Inches(0.72) * i
    add_rect(s, Inches(0.5), y, Inches(5.5), Inches(0.58), MID_BLUE)
    add_label(s, item, Inches(0.65), y + Inches(0.1), Inches(5.2), Inches(0.38),
              font_size=14, color=WHITE, bold=(item == "Runs in PARALLEL"))

# Right column — Sonnet
add_rect(s, Inches(7.2), Inches(1.55), Inches(5.8), Inches(4.9), GRAY_BG)
add_rect(s, Inches(7.2), Inches(1.55), Inches(5.8), Inches(0.55), RED)
add_label(s, "Claude Sonnet  —  Accurate + Powerful",
          Inches(7.3), Inches(1.58), Inches(5.6), Inches(0.48),
          font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
sonnet_items = ["Final Synthesizer", "ESI Score Decision",
                "Action Plan Generation", "Conversational Chat",
                "Runs AFTER all 4 complete"]
for i, item in enumerate(sonnet_items):
    y = Inches(2.25) + Inches(0.72) * i
    add_rect(s, Inches(7.35), y, Inches(5.5), Inches(0.58), MID_BLUE)
    add_label(s, item, Inches(7.5), y + Inches(0.1), Inches(5.2), Inches(0.38),
              font_size=14, color=WHITE,
              bold=(item == "Runs AFTER all 4 complete"))

# Arrow between
add_label(s, "→", Inches(6.1), Inches(3.5), Inches(1.1), Inches(0.6),
          font_size=36, bold=True, color=RED, align=PP_ALIGN.CENTER)

add_label(s, "Speed where it counts  ·  Quality where it matters",
          Inches(2.0), Inches(6.6), Inches(9.3), Inches(0.4),
          font_size=13, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Architecture Diagram
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
fill_bg(s)
slide_header(s, "System Architecture", "Six Agents, One Pipeline")
footer(s)

def box(slide, label, x, y, w=Inches(2.1), h=Inches(0.52),
        bg=GRAY_BG, fg=WHITE, fs=12, bold=False):
    add_rect(slide, x, y, w, h, bg)
    add_label(slide, label, x + Inches(0.05), y + Inches(0.06),
              w - Inches(0.1), h - Inches(0.1),
              font_size=fs, color=fg, bold=bold, align=PP_ALIGN.CENTER)

def arrow(slide, x, y, vertical=True):
    if vertical:
        add_label(slide, "▼", x, y, Inches(0.4), Inches(0.35),
                  font_size=14, color=LIGHT, align=PP_ALIGN.CENTER)
    else:
        add_label(slide, "▶", x, y, Inches(0.35), Inches(0.35),
                  font_size=14, color=LIGHT, align=PP_ALIGN.CENTER)

# Row 1 — input
box(s, "Nurse Input", Inches(5.55), Inches(1.45), bg=RED, fg=WHITE, bold=True)
arrow(s, Inches(6.35), Inches(2.02))

# Row 2 — backend
box(s, "FastAPI Backend", Inches(5.55), Inches(2.4), bg=MID_BLUE, fg=WHITE, bold=True)

# Data arrows
add_label(s, "get_patient_data() ──▶ Patient Record",
          Inches(0.4), Inches(2.52), Inches(4.9), Inches(0.3),
          font_size=10, color=LIGHT)
add_label(s, "search_protocols()  ──▶ Protocol Library",
          Inches(0.4), Inches(2.85), Inches(4.9), Inches(0.3),
          font_size=10, color=LIGHT)
arrow(s, Inches(6.35), Inches(2.97))

# Row 3 — parallel agents
add_rect(s, Inches(0.35), Inches(3.38), Inches(12.6), Inches(0.3), MONO_BG)
add_label(s, "ThreadPoolExecutor  —  4 agents in parallel",
          Inches(0.35), Inches(3.38), Inches(12.6), Inches(0.3),
          font_size=11, color=GOLD, bold=True, align=PP_ALIGN.CENTER)

agent_boxes = [
    ("Vitals\nAnalyzer", Inches(0.4)),
    ("Symptom\nClassifier", Inches(3.2)),
    ("Protocol\nMatcher", Inches(6.0)),
    ("Bed\nAllocator", Inches(8.8)),
]
for label, x in agent_boxes:
    box(s, label, x, Inches(3.72), w=Inches(2.5), h=Inches(0.9),
        bg=MID_BLUE, fg=WHITE, fs=12, bold=True)
    add_label(s, "Haiku", x + Inches(0.9), Inches(4.65), Inches(0.8), Inches(0.28),
              font_size=9, color=GOLD, align=PP_ALIGN.CENTER)

# Merge arrow
add_label(s, "▼              ▼               ▼              ▼",
          Inches(0.4), Inches(5.0), Inches(11.3), Inches(0.35),
          font_size=14, color=LIGHT)

# Row 4 — synthesizer
box(s, "Synthesizer  (Sonnet)", Inches(4.5), Inches(5.38),
    w=Inches(4.3), h=Inches(0.55), bg=RED, fg=WHITE, bold=True, fs=14)
arrow(s, Inches(6.35), Inches(5.98))

# Row 5 — report
box(s, "generate_triage_report()  →  Structured JSON",
    Inches(3.0), Inches(6.38), w=Inches(7.3), h=Inches(0.5),
    bg=GRAY_BG, fg=GOLD, bold=True, fs=13)

add_label(s, "SSE stream → React UI",
          Inches(10.5), Inches(6.42), Inches(2.5), Inches(0.4),
          font_size=11, color=LIGHT, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — The Six AI Agents
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
fill_bg(s)
slide_header(s, "The Six AI Agents", "Each Agent Has a Specialized Role")
footer(s)

headers = ["Agent", "Model", "Role"]
rows = [
    ["Vitals Analyzer",    "Haiku",  "Scores vital abnormalities 1–5, flags hemodynamic instability"],
    ["Symptom Classifier", "Haiku",  "Identifies red-flag presentations: ACS, stroke, sepsis, respiratory failure"],
    ["Protocol Matcher",   "Haiku",  "Maps presentation to evidence-based protocols and ordered interventions"],
    ["Bed Allocator",      "Haiku",  "Recommends care area (trauma bay, resus, fast track) and required equipment"],
    ["Synthesizer",        "Sonnet", "Integrates all findings into a final ESI score and action checklist"],
    ["Chat Agent",         "Haiku",  "Answers nurse follow-up questions with full report context"],
]
col_widths = [Inches(2.4), Inches(1.5), Inches(9.0)]
simple_table(s, headers, rows,
             Inches(0.35), Inches(1.6),
             col_widths, row_height=Inches(0.52), font_size=14)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Anthropic Tool Use
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
fill_bg(s)
slide_header(s, "Anthropic Tool Use", "Agents Don't Just Talk — They Call Functions")
footer(s)

add_label(s, "Claude's Tool Use feature lets agents call structured functions mid-conversation",
          Inches(0.45), Inches(1.6), Inches(12.5), Inches(0.45),
          font_size=15, color=LIGHT)

headers = ["Tool", "What It Does"]
rows = [
    ["get_patient_data",       "Fetches the full patient record by ID"],
    ["search_protocols",       "Keyword-matches clinical protocols to the chief complaint"],
    ["check_bed_availability", "Returns real-time bed counts by care area"],
    ["generate_triage_report", "Assembles and returns the structured report dict (Synthesizer only)"],
]
col_widths = [Inches(3.8), Inches(9.0)]
simple_table(s, headers, rows,
             Inches(0.35), Inches(2.1),
             col_widths, row_height=Inches(0.52), font_size=14)

# Agentic loop diagram
add_rect(s, Inches(0.35), Inches(4.55), Inches(12.6), Inches(1.85), MONO_BG)
add_label(s, "Agentic Loop",
          Inches(0.45), Inches(4.6), Inches(4.0), Inches(0.35),
          font_size=11, bold=True, color=GOLD)
loop_steps = [
    "Send message", "→", "Tool call returned", "→",
    "Dispatch function", "→", "Re-submit result", "→", "Repeat until done"
]
x_pos = Inches(0.45)
for step in loop_steps:
    is_arrow = step == "→"
    w = Inches(0.4) if is_arrow else Inches(1.85)
    add_label(s, step, x_pos, Inches(5.1), w, Inches(0.45),
              font_size=12,
              color=LIGHT if is_arrow else WHITE,
              bold=(not is_arrow),
              align=PP_ALIGN.CENTER)
    x_pos += w + Inches(0.05)

add_label(s,
          "Tool use enforces structured output — the final report is always machine-readable JSON, not LLM prose",
          Inches(0.45), Inches(6.15), Inches(12.5), Inches(0.45),
          font_size=13, italic=True, color=GOLD)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Streaming and Real-Time UX
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
fill_bg(s)
slide_header(s, "Streaming & Real-Time UX", "The Nurse Sees Progress as It Happens")
footer(s)

pts = [
    "Backend uses Server-Sent Events (SSE) — a persistent HTTP stream from server to browser",
    "As each specialist agent completes, a status event fires with the agent name and elapsed time",
    "UI displays a 6-agent grid: idle → active → done  (with per-agent completion time badge)",
    "Final complete event delivers the full report JSON",
    "No polling, no WebSockets — SSE is lightweight and browser-native",
]
bullets(s, pts, Inches(0.45), Inches(1.65), Inches(7.8), Inches(4.5), font_size=16)

# SSE event stream mockup
add_rect(s, Inches(8.7), Inches(1.5), Inches(4.3), Inches(5.3), MONO_BG)
add_label(s, "SSE Event Stream", Inches(8.8), Inches(1.55), Inches(4.1), Inches(0.35),
          font_size=11, bold=True, color=GOLD)
events = [
    ('event: status', LIGHT),
    ('data: {"agent":"vitals","done":true,"time":1.2}', WHITE),
    ('', WHITE),
    ('event: status', LIGHT),
    ('data: {"agent":"symptoms","done":true,"time":1.8}', WHITE),
    ('', WHITE),
    ('event: status', LIGHT),
    ('data: {"agent":"protocol","done":true,"time":2.1}', WHITE),
    ('', WHITE),
    ('event: complete', RGBColor(0x22,0xC5,0x5E)),
    ('data: {"report":{...}}', WHITE),
]
for i, (line, c) in enumerate(events):
    add_label(s, line, Inches(8.85), Inches(2.0) + Inches(0.38) * i,
              Inches(4.1), Inches(0.36), font_size=9, color=c)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Iterative Development
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
fill_bg(s)
slide_header(s, "Iterative Development", "Built and Improved in Six Documented Sprints")
footer(s)

headers = ["Sprint", "Key Change", "Impact"]
rows = [
    ["1", "Core pipeline, sequential agents",              "Baseline working system"],
    ["2", "Parallel agent execution",                      "~75% latency reduction"],
    ["3", "Two-tier model routing (Haiku / Sonnet)",       "Cost reduction, quality maintained"],
    ["4", "Pydantic validation + failure guards",          "Reliability, 3-failure abort"],
    ["5", "Specialist retry on timeout + thread offload",  "Stability under load"],
    ["6", "Live bed decrement, per-agent timing badges",   "Demo realism, UX polish"],
]
col_widths = [Inches(1.0), Inches(6.8), Inches(4.8)]
simple_table(s, headers, rows,
             Inches(0.35), Inches(1.6),
             col_widths, row_height=Inches(0.52), font_size=14)

add_label(s,
          "Each sprint is documented in CHANGES_N.md — mirrors real agile development",
          Inches(0.45), Inches(6.6), Inches(12.5), Inches(0.4),
          font_size=12, italic=True, color=LIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Key Engineering Decisions
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
fill_bg(s)
slide_header(s, "Key Engineering Decisions", "Design Choices and Their Trade-offs")
footer(s)

decisions = [
    ("Parallel > Sequential agents",
     "Thread management complexity vs. latency. Parallel won decisively — ~75% faster."),
    ("Two-tier model routing",
     "Haiku for data gathering (fast/cheap), Sonnet only for synthesis (accurate where it matters)."),
    ("SSE over WebSockets",
     "SSE is server-to-client only — that's all we need. Simpler, browser-native, no library."),
    ("Tool use for structured output",
     "Adds round-trip latency but ensures the final report is always machine-readable JSON."),
    ("Mock data only",
     "Removes EHR compliance complexity for this prototype. Production: replace with FHIR."),
]
for i, (title, detail) in enumerate(decisions):
    y = Inches(1.6) + Inches(1.05) * i
    add_rect(s, Inches(0.35), y, Inches(12.6), Inches(0.95), GRAY_BG)
    add_rect(s, Inches(0.35), y, Inches(0.12), Inches(0.95), RED)
    add_label(s, title, Inches(0.6), y + Inches(0.04), Inches(12.1), Inches(0.35),
              font_size=14, bold=True, color=WHITE)
    add_label(s, detail, Inches(0.6), y + Inches(0.42), Inches(12.1), Inches(0.45),
              font_size=12, color=LIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Limitations & Future Work
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
fill_bg(s)
slide_header(s, "Limitations & Future Work", "What This Prototype Does Not Do (Yet)")
footer(s)

# Left — limitations
add_rect(s, Inches(0.35), Inches(1.55), Inches(5.9), Inches(5.25), GRAY_BG)
add_rect(s, Inches(0.35), Inches(1.55), Inches(5.9), Inches(0.48), RED)
add_label(s, "Current Limitations",
          Inches(0.45), Inches(1.58), Inches(5.7), Inches(0.42),
          font_size=14, bold=True, color=WHITE)
limits = [
    "6 hardcoded mock patients — no real EHR data",
    "No authentication or role-based access control",
    "Protocol search is keyword-based, not semantic",
    "No audit log or persistence layer",
    "Not clinically validated against real outcomes",
]
bullets(s, limits, Inches(0.45), Inches(2.2), Inches(5.6), Inches(4.0), font_size=14)

# Right — future work
add_rect(s, Inches(6.45), Inches(1.55), Inches(6.5), Inches(5.25), GRAY_BG)
add_rect(s, Inches(6.45), Inches(1.55), Inches(6.5), Inches(0.48), MID_BLUE)
add_label(s, "Future Work",
          Inches(6.55), Inches(1.58), Inches(6.3), Inches(0.42),
          font_size=14, bold=True, color=GOLD)
future = [
    "FHIR/HL7 integration with real EHR systems",
    "Vector-based semantic protocol search",
    "HIPAA-compliant deployment with audit logging",
    "Outcome tracking to validate ESI accuracy",
    "Multi-patient dashboard for charge nurses",
    "Mobile-responsive UI for bedside tablets",
]
bullets(s, future, Inches(6.55), Inches(2.2), Inches(6.2), Inches(4.0), font_size=14)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Closing
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
fill_bg(s)
add_rect(s, Inches(0), Inches(0), W, Inches(0.12), RED)
add_rect(s, Inches(0), Inches(0), Inches(0.18), H, RED)
footer(s)

add_label(s, "TriageIQ Puts an AI Team Behind Every Nurse",
          Inches(0.45), Inches(0.8), Inches(12.5), Inches(1.0),
          font_size=34, bold=True, color=WHITE)

takeaways = [
    "Emergency triage is high-stakes and cognitively demanding — AI can reduce that burden",
    "A multi-agent architecture mirrors how real care teams operate: specialists in parallel, attending synthesizing",
    "Tool use + structured output makes AI recommendations reliable and consistently renderable",
    "Decision-support framing keeps the nurse in command — the AI advises, the nurse decides",
]
for i, t in enumerate(takeaways):
    y = Inches(2.0) + Inches(0.95) * i
    add_rect(s, Inches(0.45), y, Inches(12.1), Inches(0.8), GRAY_BG)
    add_rect(s, Inches(0.45), y, Inches(0.1), Inches(0.8), GOLD)
    add_label(s, f"{i+1}.  {t}", Inches(0.65), y + Inches(0.1),
              Inches(11.8), Inches(0.6), font_size=15, color=WHITE)

add_rect(s, Inches(1.5), Inches(6.1), Inches(10.3), Inches(0.72), RED)
add_label(s, '"Faster triage.  Better decisions.  The nurse still leads."',
          Inches(1.5), Inches(6.15), Inches(10.3), Inches(0.62),
          font_size=20, bold=True, italic=True, color=WHITE,
          align=PP_ALIGN.CENTER)


# ── Save ──────────────────────────────────────────────────────────────────────
out = "/Users/nicholasthomas/Downloads/triageiq/TriageIQ_Presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
